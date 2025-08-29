# app/pptx_builder.py
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from typing import List, Optional, Tuple

# ---------- Parsing ----------

def parse_markdown(md: str) -> Tuple[Optional[str], Optional[str], List[Tuple[str, List[str]]]]:
    """Return (title, subtitle, slides) where slides = [(slide_title, bullets)]
    Markdown rules:
      - First top-level heading `#` → deck title (first paragraph after it → subtitle if not a list)
      - `##` → new slide title; following paragraphs and lists → bullets
      - `-`, `*`, `+` → bullets
    """
    lines = md.replace('\r\n', '\n').split('\n')
    title = None
    subtitle = None
    slides: List[Tuple[str, List[str]]] = []

    i = 0
    # Capture deck title / subtitle
    if i < len(lines) and lines[i].startswith('# '):
        title = lines[i][2:].strip()
        i += 1
        # eat blank lines
        while i < len(lines) and not lines[i].strip():
            i += 1
        # subtitle if next non-blank is not a heading/list
        if i < len(lines) and not re.match(r'^(#|\-|\*|\+)', lines[i].strip()):
            subtitle = lines[i].strip()
            i += 1

    # Collect slides
    cur_title = None
    cur_bullets: List[str] = []

    def push_slide():
        nonlocal cur_title, cur_bullets
        if cur_title:
            slides.append((cur_title, [b for b in cur_bullets if b.strip()]))
        cur_title, cur_bullets = None, []

    while i < len(lines):
        line = lines[i]
        if line.startswith('## '):
            push_slide()
            cur_title = line[3:].strip()
        elif re.match(r'^(\-|\*|\+)\s+', line):
            cur_bullets.append(re.sub(r'^(\-|\*|\+)\s+', '', line).strip())
        elif line.strip() and not line.startswith('#'):
            # Treat paragraphs as bullets
            cur_bullets.append(line.strip())
        i += 1

    push_slide()
    return title, subtitle, slides


def parse_text(txt: str) -> Tuple[Optional[str], Optional[str], List[Tuple[str, List[str]]]]:
    """Text mode: use '---' to separate slides. Title/subtitle can be the first two non-empty lines.
    If no separators, auto-chunk by paragraph.
    """
    content = txt.replace('\r\n', '\n')
    parts = [p.strip() for p in re.split(r'\n\s*---\s*\n', content) if p.strip()]
    title, subtitle = None, None
    slides: List[Tuple[str, List[str]]] = []

    def split_title_body(block: str) -> Tuple[str, List[str]]:
        lines = [l.strip() for l in block.split('\n') if l.strip()]
        if not lines:
            return ("Slide", [])
        stitle = lines[0]
        bullets = []
        for ln in lines[1:]:
            bullets.extend([s.strip() for s in re.split(r'[•\u2022]\s*', ln) if s.strip()])
        if not bullets:
            # split into sentences as bullets
            bullets = [s.strip() for s in re.split(r'(?<=[.!?])\s+', ' '.join(lines[1:])) if s.strip()]
        return (stitle, bullets)

    if parts:
        # First block’s first lines → deck title/subtitle heuristics
        first_lines = [l for l in parts[0].split('\n') if l.strip()]
        if first_lines:
            title = first_lines[0].strip()
            if len(first_lines) > 1:
                subtitle = first_lines[1].strip()
        # each part → slide
        for block in parts:
            s_title, s_bullets = split_title_body(block)
            slides.append((s_title, s_bullets))
    else:
        # No separators; auto paragraph chunking
        paras = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]
        if paras:
            title = paras[0][:80]
        for p in paras:
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', p) if s.strip()]
            if not sentences:
                continue
            s_title = sentences[0][:80]
            slides.append((s_title, sentences[1:] or sentences[:3]))

    return title, subtitle, slides

# ---------- PPTX generation ----------

def find_layout(prs: Presentation, name_keywords: List[str], fallback_idx: int) -> int:
    name_keywords = [kw.lower() for kw in name_keywords]
    for idx, layout in enumerate(prs.slide_layouts):
        nm = (layout.name or '').lower()
        if any(kw in nm for kw in name_keywords):
            return idx
    return fallback_idx


def build_presentation(
    text: str,
    mode: str = "auto",  # auto|markdown|text
    template_path: Optional[str] = None,
    brand_rgb: Optional[Tuple[int, int, int]] = None,
) -> Presentation:
    prs = Presentation(template_path) if template_path else Presentation()

    if mode == 'markdown' or (mode == 'auto' and text.strip().startswith('# ')):
        title, subtitle, slides = parse_markdown(text)
    else:
        title, subtitle, slides = parse_text(text)

    title_layout_idx = find_layout(prs, ["title slide"], 0)
    content_layout_idx = find_layout(prs, ["title and content", "content"], 1 if len(prs.slide_layouts) > 1 else 0)

    # Title slide
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[title_layout_idx])
        if slide.shapes.title:
            slide.shapes.title.text = title
        # try subtitle placeholder (usually index 1)
        try:
            sub = slide.placeholders[1]
            if subtitle:
                sub.text = subtitle
        except Exception:
            pass

    # Content slides
    for stitle, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[content_layout_idx])
        if slide.shapes.title:
            slide.shapes.title.text = stitle or ""
        # content placeholder (usually index 1)
        body = None
        for ph in slide.placeholders:
            if ph.is_placeholder and ph.placeholder_format.type not in (0,):  # 0=TITLE
                body = ph
                break
        if not body:
            # create textbox if template lacks body placeholder
            left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
            body = slide.shapes.add_textbox(left, top, width, height)
        tf = getattr(body, 'text_frame', body)
        tf.clear()
        if bullets:
            p = tf.paragraphs[0]
            p.text = bullets[0]
            p.level = 0
            for b in bullets[1:]:
                r = tf.add_paragraph()
                r.text = b
                r.level = 0
        else:
            tf.paragraphs[0].text = ""

    # Optional: unify brand color for titles (without fighting theme)
    if brand_rgb:
        rgb = RGBColor(*brand_rgb)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes.title:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = rgb
                            run.font.bold = True
                            run.font.size = Pt(40)

    return prs